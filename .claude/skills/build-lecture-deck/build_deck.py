"""Convert a curriculum lecture slide-outline (``.md``) into a PowerPoint (``.pptx``) deck.

Input is one lecture slide-outline written per ``docs/origin/LECTURES.md``:

    # <deck title>
    ```yaml
    title: ...
    keywords: ...
    estimated duration: NN
    ```
    > **Lesson:** ...            # metadata blockquote (optional)
    ## section N. <section title>
    ### slide N.M <slide title>
    - a bullet
      - a sub-bullet
    | col | col |                # a GitHub-style markdown table
    | --- | --- |
    | a   | b   |
    [↑ Back to top](#...)        # dropped

Slide mapping:
    H1 + yaml + blockquote   -> one title slide
    ``## section N. Title``  -> one section-divider slide
    ``### slide N.M Title``   -> one content slide (bullets + tables)
    ``diagram:`` / ``table:`` lead bullets -> speaker notes (authoring directives, not slide text)

Usage:
    uv run python .claude/skills/build-lecture-deck/build_deck.py <lecture.md> [--out deck.pptx]

python-pptx ships incomplete type information, so the render helpers below annotate its
presentation / slide / paragraph objects as ``Any`` on purpose (this file is outside the
repo's pyright ``include``; the geometry helpers still use the real ``Length`` type).
"""

from __future__ import annotations

import argparse
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal

from pptx import Presentation
from pptx.util import Inches, Length, Pt

# --- parsing -----------------------------------------------------------------

_H1 = re.compile(r"^#\s+(.+?)\s*$")
_SECTION = re.compile(r"^##\s+section\s+([0-9]+)\.\s*(.+?)\s*$", re.IGNORECASE)
_SLIDE = re.compile(r"^###\s+slide\s+([0-9]+\.[0-9]+)\s+(.+?)\s*$", re.IGNORECASE)
_BULLET = re.compile(r"^(?P<indent>\s*)(?:[-*]|[0-9]+\.)\s+(?P<text>.+?)\s*$")
_BACK_TO_TOP = re.compile(r"^\s*\[.\s*Back to top\]", re.IGNORECASE)
_LINK = re.compile(r"\[([^\]]+)\]\([^)]+\)")
_INLINE = re.compile(r"(\*\*.+?\*\*|\*[^*\n]+?\*|`[^`]+`)")


@dataclass
class Bullet:
    text: str
    level: int


@dataclass
class Table:
    header: list[str]
    rows: list[list[str]]


Block = Bullet | Table


@dataclass
class SlideSpec:
    kind: Literal["title", "section", "content"]
    title: str
    number: str = ""
    subtitle: str = ""
    blocks: list[Block] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)


def _clean(text: str) -> str:
    """Strip markdown links to their visible text and trim whitespace."""
    return _LINK.sub(r"\1", text).strip()


def _split_table_row(line: str) -> list[str]:
    return [c.strip() for c in line.strip().strip("|").split("|")]


def _title_slide(deck_title: str, meta: dict[str, str], blockquote: list[str]) -> SlideSpec:
    """Build the opening title slide from the H1, yaml block, and metadata blockquote."""
    bq = _clean(" ".join(blockquote)).replace("**", "")
    parts: list[str] = []
    lesson = re.search(r"Lesson:\s*([A-Za-z0-9]+)", bq)
    if lesson:
        parts.append(lesson.group(1))
    duration = meta.get("estimated duration")
    if duration:
        parts.append(f"~{duration} min")
    anchor = re.search(r"Anchor model[^:]*:\s*([^.·]+)", bq)
    if anchor:
        parts.append(anchor.group(1).strip())
    return SlideSpec(
        kind="title",
        title=deck_title,
        subtitle="  ·  ".join(parts),
        notes=[bq] if bq else [],
    )


def parse_outline(md: str) -> list[SlideSpec]:
    """Parse a lecture slide-outline into an ordered list of slide specs."""
    lines = md.splitlines()
    n = len(lines)
    slides: list[SlideSpec] = []

    # --- header region: H1, then optional yaml fence, then leading blockquote ---
    deck_title = ""
    meta: dict[str, str] = {}
    blockquote: list[str] = []

    i = 0
    while i < n:
        h1 = _H1.match(lines[i])
        i += 1
        if h1:
            deck_title = _clean(h1.group(1))
            break

    while i < n and lines[i].strip() == "":
        i += 1
    if i < n and lines[i].strip().startswith("```yaml"):
        i += 1
        while i < n and not lines[i].strip().startswith("```"):
            key, _, val = lines[i].partition(":")
            if val:
                meta[key.strip()] = val.strip().strip('"')
            i += 1
        i += 1  # consume the closing fence

    while i < n and not _SECTION.match(lines[i]) and not _SLIDE.match(lines[i]):
        if lines[i].lstrip().startswith(">"):
            blockquote.append(lines[i].lstrip()[1:].strip())
        i += 1

    slides.append(_title_slide(deck_title, meta, blockquote))

    # --- body: sections and slides ---
    current: SlideSpec | None = None
    pending_table: list[str] = []
    last_bullet: Bullet | None = None

    def flush_table() -> None:
        nonlocal pending_table, last_bullet
        if pending_table and current is not None and len(pending_table) >= 2:
            header = [_clean(c) for c in _split_table_row(pending_table[0])]
            rows = [[_clean(c) for c in _split_table_row(r)] for r in pending_table[2:]]
            current.blocks.append(Table(header=header, rows=rows))
        pending_table = []
        last_bullet = None

    for line in lines[i:]:
        section = _SECTION.match(line)
        slide = _SLIDE.match(line)
        if section:
            flush_table()
            current = SlideSpec(
                kind="section", title=_clean(section.group(2)), number=section.group(1)
            )
            slides.append(current)
            last_bullet = None
            continue
        if slide:
            flush_table()
            current = SlideSpec(kind="content", title=_clean(slide.group(2)), number=slide.group(1))
            slides.append(current)
            last_bullet = None
            continue
        if current is None or _BACK_TO_TOP.match(line):
            continue

        if line.strip().startswith("|"):
            pending_table.append(line)
            continue
        if pending_table:
            flush_table()

        if line.strip() == "":
            last_bullet = None
            continue

        bullet = _BULLET.match(line)
        if bullet:
            text = _clean(bullet.group("text"))
            if text.lower().startswith(("diagram:", "table:")):
                current.notes.append(text)
                last_bullet = None
                continue
            item = Bullet(text=text, level=min(len(bullet.group("indent")) // 2, 4))
            current.blocks.append(item)
            last_bullet = item
        elif last_bullet is not None:
            # a wrapped continuation line of the previous bullet
            last_bullet.text = f"{last_bullet.text} {_clean(line.strip())}"

    flush_table()
    return slides


# --- rendering ---------------------------------------------------------------


def _add_inline(paragraph: Any, text: str) -> None:
    """Append runs to ``paragraph``, honoring ``**bold**``, ``*italic*``, and ``code``."""
    for token in _INLINE.split(text):
        if not token:
            continue
        run = paragraph.add_run()
        if token.startswith("**") and token.endswith("**"):
            run.text = token[2:-2]
            run.font.bold = True
        elif token.startswith("`") and token.endswith("`"):
            run.text = token[1:-1]
            run.font.name = "Consolas"
        elif token.startswith("*") and token.endswith("*"):
            run.text = token[1:-1]
            run.font.italic = True
        else:
            run.text = token


def _set_notes(slide: Any, notes: list[str]) -> None:
    if notes:
        slide.notes_slide.notes_text_frame.text = "\n".join(notes)


def _add_bullets(slide: Any, bullets: list[Bullet], top: Length, height: Length) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), height)
    frame = box.text_frame
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.level = bullet.level
        para.space_after = Pt(6)
        marker = "•  " if bullet.level == 0 else "◦  "
        _add_inline(para, marker + bullet.text)
        for run in para.runs:
            run.font.size = Pt(16 - 2 * min(bullet.level, 2))


def _add_table(slide: Any, table: Table, top: Length) -> Length:
    n_cols = len(table.header)
    for row in table.rows:
        n_cols = max(n_cols, len(row))
    n_rows = len(table.rows) + 1
    height = Inches(0.4 * n_rows)
    grid = slide.shapes.add_table(n_rows, n_cols, Inches(0.6), top, Inches(12.1), height).table

    for col, text in enumerate(table.header):
        para = grid.cell(0, col).text_frame.paragraphs[0]
        _add_inline(para, text)
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(11)

    for r, row in enumerate(table.rows, start=1):
        for col in range(n_cols):
            para = grid.cell(r, col).text_frame.paragraphs[0]
            _add_inline(para, row[col] if col < len(row) else "")
            for run in para.runs:
                run.font.size = Pt(10)

    return Inches(top.inches + height.inches + 0.2)


def _render_content(prs: Any, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = f"{spec.number}  {spec.title}".strip()

    bullets = [b for b in spec.blocks if isinstance(b, Bullet)]
    tables = [b for b in spec.blocks if isinstance(b, Table)]

    if bullets:
        _add_bullets(slide, bullets, Inches(1.5), Inches(1.9) if tables else Inches(5.4))
    top: Length = Inches(3.5) if bullets else Inches(1.7)
    for table in tables:
        top = _add_table(slide, table, top)
    _set_notes(slide, spec.notes)


def _render_section(prs: Any, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section Header
    slide.shapes.title.text = f"{spec.number}. {spec.title}"
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx != 0:
            placeholder.text = "Section"
            break
    _set_notes(slide, spec.notes)


def _render_title(prs: Any, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide.shapes.title.text = spec.title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = spec.subtitle
    _set_notes(slide, spec.notes)


def build_deck(slides: list[SlideSpec], out_path: Path) -> None:
    """Render parsed slide specs to a 16:9 ``.pptx`` at ``out_path``."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for spec in slides:
        if spec.kind == "title":
            _render_title(prs, spec)
        elif spec.kind == "section":
            _render_section(prs, spec)
        else:
            _render_content(prs, spec)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Convert a curriculum lecture slide-outline (.md) into a .pptx deck."
    )
    parser.add_argument("input", type=Path, help="Path to the lecture slide-outline .md file.")
    parser.add_argument("--out", type=Path, default=None, help="Output .pptx path.")
    args = parser.parse_args()

    slides = parse_outline(args.input.read_text(encoding="utf-8"))
    out = args.out or args.input.with_suffix(".pptx")
    build_deck(slides, out)

    sections = sum(1 for s in slides if s.kind == "section")
    content = sum(1 for s in slides if s.kind == "content")
    print(f"Wrote {out} — {len(slides)} slides ({sections} sections, {content} content).")


if __name__ == "__main__":
    main()
